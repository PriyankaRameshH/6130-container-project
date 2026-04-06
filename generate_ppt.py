#!/usr/bin/env python3
"""Generate a 10-slide master's-level presentation for the eBPF Container Security project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme Colors ──────────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
ACCENT_BLUE   = RGBColor(0x00, 0xB4, 0xD8)   # Cyan accent
ACCENT_GREEN  = RGBColor(0x00, 0xF5, 0xD4)   # Teal green
ACCENT_RED    = RGBColor(0xFF, 0x47, 0x6B)    # Alert red
ACCENT_ORANGE = RGBColor(0xFF, 0xA6, 0x2B)    # Warning orange
ACCENT_PURPLE = RGBColor(0xB5, 0x83, 0xFF)    # Purple accent
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xBE, 0xC5)
MED_GRAY      = RGBColor(0x78, 0x90, 0x9C)
DARK_CARD     = RGBColor(0x1A, 0x23, 0x3B)    # Card background
CODE_BG       = RGBColor(0x12, 0x1B, 0x2E)    # Code block bg

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_dark_bg(slide):
    """Fill slide with dark background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_BLUE):
    """Add multi-line bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(2)
        # Bullet character
        run_b = p.add_run()
        run_b.text = "◆  "
        run_b.font.size = Pt(font_size - 2)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = "Segoe UI"
        # Text
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = "Segoe UI"
    return txBox


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    """Add a styled code block."""
    shape = add_shape(slide, left, top, width, height, CODE_BG, ACCENT_BLUE)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = ACCENT_GREEN
        run.font.name = "Consolas"
    return shape


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a horizontal accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total=10):
    """Add slide number at bottom-right."""
    add_text_box(slide, Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4),
                 f"{num} / {total}", font_size=11, color=MED_GRAY,
                 alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title Slide
# ══════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide1)

# Decorative top accent bar
add_shape(slide1, Inches(0), Inches(0), W, Pt(4), ACCENT_BLUE)

# Central content area
add_text_box(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "CONTAINER RUNTIME SECURITY", font_size=44, color=WHITE, bold=True)
add_accent_line(slide1, Inches(1), Inches(2.5), Inches(3), ACCENT_BLUE)
add_text_box(slide1, Inches(1), Inches(2.7), Inches(11), Inches(1),
             "Real-Time Detection of Container Escape Attacks Using eBPF",
             font_size=26, color=ACCENT_BLUE, bold=False)

add_text_box(slide1, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
             "An eBPF-Based Kernel Tracing Approach to Monitor and Alert on "
             "Five Distinct Container Escape Techniques in Docker Environments",
             font_size=16, color=LIGHT_GRAY)

# Tech stack badges
badges = ["eBPF/libbpf", "C (Kernel + Userspace)", "Docker", "Linux Kernel 6.x", "CO-RE"]
x_start = Inches(1)
for badge in badges:
    bw = Inches(len(badge) * 0.12 + 0.4)
    shape = add_shape(slide1, x_start, Inches(5.3), bw, Inches(0.4), DARK_CARD, ACCENT_BLUE)
    tf = shape.text_frame
    tf.margin_left = Pt(8)
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = badge
    run.font.size = Pt(12)
    run.font.color.rgb = ACCENT_BLUE
    run.font.name = "Consolas"
    p.alignment = PP_ALIGN.CENTER
    x_start += bw + Inches(0.15)

add_text_box(slide1, Inches(1), Inches(6.2), Inches(5), Inches(0.5),
             "Course: CS 6130  |  Platform: Ubuntu 24.04 / aarch64",
             font_size=13, color=MED_GRAY)
add_slide_number(slide1, 1)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 2 — Problem Statement & Motivation
# ══════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide2)
add_shape(slide2, Inches(0), Inches(0), W, Pt(4), ACCENT_RED)

add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "PROBLEM STATEMENT & MOTIVATION", font_size=32, color=WHITE, bold=True)
add_accent_line(slide2, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_RED)

# Left column - The Problem
add_shape(slide2, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2), DARK_CARD, ACCENT_RED)
add_text_box(slide2, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
             "⚠  The Problem", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide2, Inches(0.8), Inches(2.2), Inches(5.2), Inches(4.2), [
    "Containers share the host kernel — a single\nvulnerability can compromise the entire host",
    "Docker misconfigurations (--privileged, socket\nmounts) create direct escape paths",
    "CVEs like CVE-2022-0492 allow cgroup-based\nhost code execution from inside containers",
    "Traditional security tools operate at user-space\nlevel and miss kernel-level escape attempts",
    "Short-lived container processes exit before\nlog-based detectors can even inspect them",
], font_size=14, bullet_color=ACCENT_RED)

# Right column - Our Solution
add_shape(slide2, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), DARK_CARD, ACCENT_GREEN)
add_text_box(slide2, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.5),
             "✓  Our Approach", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide2, Inches(7.1), Inches(2.2), Inches(5.2), Inches(4.2), [
    "eBPF tracepoints hook directly into kernel\nsyscall entry — zero-miss detection",
    "BPF CO-RE captures process metadata (mntns,\ncaps) at syscall time — defeats race conditions",
    "Ring buffer delivers events to user-space\ndetector for enrichment and classification",
    "5 specialized detection functions cover all\nmajor container escape categories",
    "Real attacks against live Docker containers\nvalidate every detection rule end-to-end",
], font_size=14, bullet_color=ACCENT_GREEN)

add_slide_number(slide2, 2)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 3 — System Architecture
# ══════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide3)
add_shape(slide3, Inches(0), Inches(0), W, Pt(4), ACCENT_BLUE)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "SYSTEM ARCHITECTURE", font_size=32, color=WHITE, bold=True)
add_accent_line(slide3, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE)

# Kernel Space Box
add_shape(slide3, Inches(0.5), Inches(1.6), Inches(4), Inches(5.2), DARK_CARD, ACCENT_ORANGE)
add_text_box(slide3, Inches(0.7), Inches(1.7), Inches(3.6), Inches(0.4),
             "KERNEL SPACE  (eBPF)", font_size=18, color=ACCENT_ORANGE, bold=True)

tracepoints = [
    ("sys_enter_mount", "mount() syscalls"),
    ("sys_enter_setns", "setns() namespace joins"),
    ("sys_enter_openat", "openat() file access"),
    ("sys_enter_connect", "connect() socket calls"),
]
y = Inches(2.3)
for tp_name, tp_desc in tracepoints:
    add_shape(slide3, Inches(0.8), y, Inches(3.4), Inches(0.55), CODE_BG, ACCENT_ORANGE)
    txBox = slide3.shapes.add_textbox(Inches(1.0), y + Pt(2), Inches(3.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = tp_name
    r1.font.size = Pt(12)
    r1.font.color.rgb = ACCENT_GREEN
    r1.font.name = "Consolas"
    r1.font.bold = True
    r2 = p.add_run()
    r2.text = f"  →  {tp_desc}"
    r2.font.size = Pt(11)
    r2.font.color.rgb = LIGHT_GRAY
    r2.font.name = "Segoe UI"
    y += Inches(0.7)

# CO-RE + Maps
add_shape(slide3, Inches(0.8), y + Inches(0.15), Inches(3.4), Inches(1.1), CODE_BG, ACCENT_BLUE)
txBox = slide3.shapes.add_textbox(Inches(1.0), y + Inches(0.2), Inches(3.0), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "BPF CO-RE Reads:"
r.font.size = Pt(12); r.font.color.rgb = ACCENT_BLUE; r.font.bold = True; r.font.name = "Segoe UI"
for item in ["• task→nsproxy→mnt_ns→ns.inum", "• task→cred→cap_effective.val"]:
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = item
    r2.font.size = Pt(11); r2.font.color.rgb = ACCENT_GREEN; r2.font.name = "Consolas"

# Arrow
add_text_box(slide3, Inches(4.5), Inches(3.3), Inches(0.8), Inches(1),
             "▶\n▶\n▶", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Ring Buffer
add_shape(slide3, Inches(5.0), Inches(2.8), Inches(2.5), Inches(1.6), DARK_CARD, ACCENT_BLUE)
add_text_box(slide3, Inches(5.2), Inches(2.9), Inches(2.1), Inches(0.4),
             "RING BUFFER", font_size=16, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide3, Inches(5.2), Inches(3.35), Inches(2.1), Inches(0.9),
             "BPF_MAP_TYPE_RINGBUF\n16 MB capacity\n250ms poll interval", font_size=11,
             color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Arrow
add_text_box(slide3, Inches(7.5), Inches(3.3), Inches(0.8), Inches(1),
             "▶\n▶\n▶", font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# User Space Box
add_shape(slide3, Inches(8.0), Inches(1.6), Inches(4.8), Inches(5.2), DARK_CARD, ACCENT_GREEN)
add_text_box(slide3, Inches(8.2), Inches(1.7), Inches(4.4), Inches(0.4),
             "USER SPACE  (C Detector)", font_size=18, color=ACCENT_GREEN, bold=True)

us_items = [
    ("1. Filter Self", "Skip detector's own PID"),
    ("2. Mntns Check", "BPF mntns ≠ host → container"),
    ("3. Enrich", "/proc/<pid>/cgroup → ID, runtime"),
    ("4. Noise Filter", "Skip dockerd, runc, systemd"),
    ("5. Detect", "5 exclusive detection functions"),
    ("6. Alert", "CRITICAL / HIGH / MEDIUM output"),
]
y = Inches(2.3)
for step, desc in us_items:
    add_shape(slide3, Inches(8.3), y, Inches(4.2), Inches(0.5), CODE_BG, ACCENT_GREEN)
    txBox = slide3.shapes.add_textbox(Inches(8.5), y + Pt(2), Inches(3.8), Inches(0.45))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = step + "  "
    r1.font.size = Pt(12); r1.font.color.rgb = ACCENT_GREEN; r1.font.bold = True; r1.font.name = "Segoe UI"
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(11); r2.font.color.rgb = LIGHT_GRAY; r2.font.name = "Segoe UI"
    y += Inches(0.6)

add_slide_number(slide3, 3)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 4 — eBPF Kernel Program (BPF Code Deep Dive)
# ══════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide4)
add_shape(slide4, Inches(0), Inches(0), W, Pt(4), ACCENT_ORANGE)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "eBPF KERNEL PROGRAM — escape_detector.bpf.c", font_size=32, color=WHITE, bold=True)
add_accent_line(slide4, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_ORANGE)
add_text_box(slide4, Inches(0.8), Inches(1.2), Inches(10), Inches(0.4),
             "191 lines  |  4 tracepoints  |  2 BPF maps  |  CO-RE portable",
             font_size=14, color=MED_GRAY)

# Event struct code block
add_code_block(slide4, Inches(0.5), Inches(1.8), Inches(5.5), Inches(4.8),
    "struct event {\n"
    "  uint64_t ts_ns;       // kernel timestamp\n"
    "  uint32_t pid, tgid;   // thread & process ID\n"
    "  uint32_t uid, gid;    // user & group ID\n"
    "  uint32_t mntns;       // mount namespace\n"
    "  uint32_t event_type;  // MOUNT|OPENAT|CONNECT|SETNS\n"
    "  int32_t  fd;          // file descriptor\n"
    "  uint32_t flags;       // syscall flags\n"
    "  uint32_t family;      // socket family (AF_UNIX)\n"
    "  uint64_t cgroup_id;   // cgroup identifier\n"
    "  uint64_t cap_eff;     // effective capabilities\n"
    "  char comm[16];        // process name\n"
    "  char path[128];       // file/socket path\n"
    "  char extra[128];      // mount target/extra\n"
    "};", font_size=11)

# Right side — Key concepts
add_shape(slide4, Inches(6.3), Inches(1.8), Inches(6.5), Inches(2.0), DARK_CARD, ACCENT_ORANGE)
add_text_box(slide4, Inches(6.6), Inches(1.9), Inches(6.0), Inches(0.4),
             "CO-RE: Compile Once, Run Everywhere", font_size=16, color=ACCENT_ORANGE, bold=True)
add_code_block(slide4, Inches(6.5), Inches(2.4), Inches(6.1), Inches(1.2),
    "struct task_struct *task = bpf_get_current_task();\n"
    "e->mntns  = BPF_CORE_READ(task, nsproxy, mnt_ns, ns.inum);\n"
    "e->cap_eff = BPF_CORE_READ(task, cred, cap_effective.val);",
    font_size=11)

# Self-filtering
add_shape(slide4, Inches(6.3), Inches(4.1), Inches(6.5), Inches(1.5), DARK_CARD, ACCENT_BLUE)
add_text_box(slide4, Inches(6.6), Inches(4.2), Inches(6.0), Inches(0.4),
             "Self-Filtering (Prevents Feedback Loop)", font_size=16, color=ACCENT_BLUE, bold=True)
add_code_block(slide4, Inches(6.5), Inches(4.65), Inches(6.1), Inches(0.8),
    "static int is_self(void) {\n"
    "  __u32 *pid = bpf_map_lookup_elem(&detector_pid, &key);\n"
    "  return (pid && *pid == (bpf_get_current_pid_tgid()>>32));\n"
    "}", font_size=11)

# Maps description
add_shape(slide4, Inches(6.3), Inches(5.9), Inches(3.0), Inches(0.85), DARK_CARD, ACCENT_GREEN)
add_text_box(slide4, Inches(6.5), Inches(5.95), Inches(2.7), Inches(0.7),
             "events  (RINGBUF 16MB)\nKernel → Userspace event pipe",
             font_size=12, color=ACCENT_GREEN, font_name="Consolas")

add_shape(slide4, Inches(9.6), Inches(5.9), Inches(3.2), Inches(0.85), DARK_CARD, ACCENT_GREEN)
add_text_box(slide4, Inches(9.8), Inches(5.95), Inches(2.9), Inches(0.7),
             "detector_pid  (ARRAY[1])\nStores detector PID for filtering",
             font_size=12, color=ACCENT_GREEN, font_name="Consolas")

add_slide_number(slide4, 4)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 5 — User-Space Detector Engine
# ══════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide5)
add_shape(slide5, Inches(0), Inches(0), W, Pt(4), ACCENT_GREEN)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "USER-SPACE DETECTOR ENGINE — detector.c", font_size=32, color=WHITE, bold=True)
add_accent_line(slide5, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_GREEN)
add_text_box(slide5, Inches(0.8), Inches(1.2), Inches(10), Inches(0.4),
             "658 lines  |  5 detection functions  |  Metadata enrichment  |  Noise filtering",
             font_size=14, color=MED_GRAY)

# Detection pipeline
detections = [
    ("DOCKER-SOCKET-ESCAPE", "CRITICAL", "connect(AF_UNIX) to /var/run/docker.sock",
     "Container accesses Docker daemon API", ACCENT_RED),
    ("PRIVILEGED-ESCAPE", "CRITICAL", "mount() + openat(/proc/1/root/*) + CAP_SYS_ADMIN",
     "Privileged container escapes via mount or /proc", ACCENT_RED),
    ("CGROUP-ESCAPE", "CRITICAL", "mount(\"cgroup\") + openat(release_agent)",
     "CVE-2022-0492 cgroup release_agent abuse", ACCENT_RED),
    ("SENSITIVE-FS-ACCESS", "HIGH", "openat(/hostfs/etc/shadow, /proc/1/environ)",
     "Container reads host credentials via bind mount", ACCENT_ORANGE),
    ("NAMESPACE-ESCAPE", "CRITICAL", "setns() + openat(/proc/1/ns/*)",
     "Container joins host namespaces via nsenter", ACCENT_RED),
]

y = Inches(1.7)
for name, sev, trigger, desc, color in detections:
    add_shape(slide5, Inches(0.5), y, Inches(12.3), Inches(0.85), DARK_CARD, color)
    # Severity badge
    sw = Inches(1.3)
    add_shape(slide5, Inches(0.7), y + Pt(6), sw, Inches(0.35), color)
    txBox = slide5.shapes.add_textbox(Inches(0.7), y + Pt(6), sw, Inches(0.35))
    tf = txBox.text_frame
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = sev
    r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = "Segoe UI"
    # Name
    add_text_box(slide5, Inches(2.2), y + Pt(2), Inches(3.5), Inches(0.35),
                 name, font_size=15, color=color, bold=True, font_name="Consolas")
    # Trigger
    add_text_box(slide5, Inches(2.2), y + Pt(22), Inches(10), Inches(0.35),
                 trigger, font_size=11, color=LIGHT_GRAY, font_name="Consolas")
    # Description
    add_text_box(slide5, Inches(6.5), y + Pt(2), Inches(6), Inches(0.35),
                 desc, font_size=13, color=WHITE)
    y += Inches(0.95)

# Exclusive ordering note
add_shape(slide5, Inches(0.5), y + Inches(0.2), Inches(12.3), Inches(0.5), DARK_CARD, ACCENT_PURPLE)
add_text_box(slide5, Inches(0.8), y + Inches(0.25), Inches(11.8), Inches(0.4),
             "⚡ Exclusive Order: Docker Socket → Cgroup → Namespace → Privileged → Sensitive FS  "
             "(first match wins — prevents cross-triggering between detectors)",
             font_size=13, color=ACCENT_PURPLE)

add_slide_number(slide5, 5)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 6 — Five Attack Scenarios
# ══════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide6)
add_shape(slide6, Inches(0), Inches(0), W, Pt(4), ACCENT_RED)

add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "FIVE REAL CONTAINER ESCAPE ATTACKS", font_size=32, color=WHITE, bold=True)
add_accent_line(slide6, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_RED)

attacks = [
    ("1", "Docker Socket\nEscape", "-v /var/run/docker.sock\n:/var/run/docker.sock",
     "curl → Docker API →\nspawn privileged container", ACCENT_RED),
    ("2", "Privileged\nEscape", "--privileged\n--pid=host",
     "mount proc → read\n/proc/1/root/etc/shadow", ACCENT_ORANGE),
    ("3", "Cgroup Escape\n(CVE-2022-0492)", "--privileged",
     "mount cgroup → write\nrelease_agent", ACCENT_ORANGE),
    ("4", "Sensitive FS\nAccess", "-v /:/hostfs:ro\n--pid=host",
     "read /hostfs/etc/shadow\n/proc/1/environ", ACCENT_BLUE),
    ("5", "Namespace\nEscape", "--privileged\n--pid=host",
     "nsenter -t 1 -m -u -n\n-i -p → full host access", ACCENT_PURPLE),
]

x = Inches(0.3)
for num, name, flags, technique, color in attacks:
    card_w = Inches(2.4)
    add_shape(slide6, x, Inches(1.5), card_w, Inches(5.3), DARK_CARD, color)
    # Number circle
    circle = slide6.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), Inches(1.7), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(24); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = "Segoe UI"

    add_text_box(slide6, x + Inches(0.15), Inches(2.55), Inches(2.1), Inches(0.7),
                 name, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Docker flags
    add_shape(slide6, x + Inches(0.1), Inches(3.35), Inches(2.2), Inches(0.9), CODE_BG)
    add_text_box(slide6, x + Inches(0.2), Inches(3.4), Inches(2.0), Inches(0.8),
                 flags, font_size=10, color=ACCENT_GREEN, font_name="Consolas", alignment=PP_ALIGN.CENTER)

    # Technique
    add_text_box(slide6, x + Inches(0.15), Inches(4.5), Inches(2.1), Inches(0.9),
                 "Technique:", font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x + Inches(0.15), Inches(4.85), Inches(2.1), Inches(0.9),
                 technique, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Docker run line
    add_shape(slide6, x + Inches(0.1), Inches(5.85), Inches(2.2), Inches(0.35), color)
    sev_text = "CRITICAL" if num != "4" else "HIGH"
    add_text_box(slide6, x + Inches(0.1), Inches(5.87), Inches(2.2), Inches(0.3),
                 sev_text, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    x += Inches(2.55)

add_text_box(slide6, Inches(0.5), Inches(6.85), Inches(12), Inches(0.4),
             "All attacks run as real Docker containers against live Docker daemon — not simulated or mocked",
             font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_slide_number(slide6, 6)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 7 — Live Detection Results
# ══════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide7)
add_shape(slide7, Inches(0), Inches(0), W, Pt(4), ACCENT_GREEN)

add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "LIVE DETECTION RESULTS", font_size=32, color=WHITE, bold=True)
add_accent_line(slide7, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_GREEN)

# Simulated terminal output
alerts = [
    ("[CRITICAL] [DOCKER-SOCKET-ESCAPE] rule=docker-socket-connect",
     "  pid=165727 comm=curl path=/var/run/docker.sock container=8c8bb2b1baee", ACCENT_RED),
    ("[CRITICAL] [PRIVILEGED-ESCAPE] rule=container-mount",
     "  pid=166052 comm=mount path=proc container=b553d4c60a5d", ACCENT_RED),
    ("[CRITICAL] [PRIVILEGED-ESCAPE] rule=proc-root-access",
     "  pid=166062 comm=cat path=/proc/1/root/etc/hostname container=b553d4c60a5d", ACCENT_RED),
    ("[CRITICAL] [CGROUP-ESCAPE] rule=cgroup-mount",
     "  pid=166366 comm=mount path=cgroup container=203b535e89c2", ACCENT_RED),
    ("[HIGH] [SENSITIVE-FS-ACCESS] rule=hostfs-credential-read",
     "  pid=167060 comm=bash path=/hostfs/etc/shadow container=7d184871bd8c", ACCENT_ORANGE),
    ("[HIGH] [SENSITIVE-FS-ACCESS] rule=host-environ-read",
     "  pid=167065 comm=bash path=/proc/1/environ container=7d184871bd8c", ACCENT_ORANGE),
    ("[CRITICAL] [NAMESPACE-ESCAPE] rule=namespace-setns",
     "  pid=167345 comm=nsenter container=80f95cc29cf1", ACCENT_RED),
    ("[HIGH] [NAMESPACE-ESCAPE] rule=namespace-probe",
     "  pid=167351 comm=nsenter path=/proc/1/ns/mnt container=80f95cc29cf1", ACCENT_ORANGE),
]

add_shape(slide7, Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.5), CODE_BG, ACCENT_GREEN)
add_text_box(slide7, Inches(0.5), Inches(1.55), Inches(3), Inches(0.4),
             "  $ sudo ./bin/detector", font_size=13, color=ACCENT_GREEN, font_name="Consolas", bold=True)

y = Inches(2.0)
for line1, line2, color in alerts:
    txBox = slide7.shapes.add_textbox(Inches(0.6), y, Inches(12), Inches(0.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = line1
    r1.font.size = Pt(11); r1.font.color.rgb = color; r1.font.name = "Consolas"; r1.font.bold = True
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = line2
    r2.font.size = Pt(10); r2.font.color.rgb = LIGHT_GRAY; r2.font.name = "Consolas"
    y += Inches(0.6)

add_slide_number(slide7, 7)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 8 — Key Challenges & Solutions
# ══════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide8)
add_shape(slide8, Inches(0), Inches(0), W, Pt(4), ACCENT_PURPLE)

add_text_box(slide8, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "KEY CHALLENGES & SOLUTIONS", font_size=32, color=WHITE, bold=True)
add_accent_line(slide8, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_PURPLE)

challenges = [
    ("Feedback Loop", "Detector reads /proc → triggers eBPF → infinite recursion",
     "BPF-side detector_pid map + userspace self_tgid filter", ACCENT_RED, ACCENT_GREEN),
    ("Short-Lived Process Race", "Container process exits before /proc/<pid>/cgroup is readable",
     "CO-RE reads mntns + cap_effective at kernel syscall time", ACCENT_RED, ACCENT_GREEN),
    ("False Positives", "systemd-oomd, dockerd, runc generate noise events",
     "Host noise comm[] allowlist + strict container-only filtering", ACCENT_ORANGE, ACCENT_GREEN),
    ("Cross-Triggering", "Attack scripts trigger multiple detection categories",
     "Exclusive detection order (first match wins) + per-attack script isolation", ACCENT_ORANGE, ACCENT_GREEN),
    ("BPF Verifier Rejection", "connect() probe failed with uninitialized stack memory",
     "__builtin_memset zeroes path/extra before bpf_probe_read_user_str()", ACCENT_ORANGE, ACCENT_GREEN),
]

y = Inches(1.6)
for title, problem, solution, prob_color, sol_color in challenges:
    # Card
    add_shape(slide8, Inches(0.5), y, Inches(12.3), Inches(0.95), DARK_CARD, ACCENT_PURPLE)
    # Title
    add_text_box(slide8, Inches(0.8), y + Pt(3), Inches(2.5), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
    # Problem
    txBox = slide8.shapes.add_textbox(Inches(3.3), y + Pt(3), Inches(4.5), Inches(0.4))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "✗  "
    r1.font.size = Pt(12); r1.font.color.rgb = prob_color; r1.font.name = "Segoe UI"; r1.font.bold = True
    r2 = p.add_run()
    r2.text = problem
    r2.font.size = Pt(11); r2.font.color.rgb = LIGHT_GRAY; r2.font.name = "Segoe UI"
    # Solution
    txBox2 = slide8.shapes.add_textbox(Inches(7.8), y + Pt(3), Inches(5), Inches(0.4))
    tf2 = txBox2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r3 = p2.add_run()
    r3.text = "✓  "
    r3.font.size = Pt(12); r3.font.color.rgb = sol_color; r3.font.name = "Segoe UI"; r3.font.bold = True
    r4 = p2.add_run()
    r4.text = solution
    r4.font.size = Pt(11); r4.font.color.rgb = LIGHT_GRAY; r4.font.name = "Segoe UI"
    y += Inches(1.05)

add_slide_number(slide8, 8)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 9 — Project Structure & Build
# ══════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide9)
add_shape(slide9, Inches(0), Inches(0), W, Pt(4), ACCENT_BLUE)

add_text_box(slide9, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "PROJECT STRUCTURE & BUILD SYSTEM", font_size=32, color=WHITE, bold=True)
add_accent_line(slide9, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE)

# Directory tree
add_code_block(slide9, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.2),
    "container-escape-detector/\n"
    "├── internal/bpf/\n"
    "│   ├── escape_detector.bpf.c   # eBPF kernel program\n"
    "│   └── escape_detector.bpf.o   # Compiled BPF object\n"
    "├── cmd/detector/\n"
    "│   └── detector.c              # User-space engine\n"
    "├── attacks/\n"
    "│   ├── attack1_docker_socket.sh\n"
    "│   ├── attack2_privileged_escape.sh\n"
    "│   ├── attack3_cgroup_escape.sh\n"
    "│   ├── attack4_sensitive_fs_access.sh\n"
    "│   ├── attack5_namespace_escape.sh\n"
    "│   └── run_real_attacks.sh      # Orchestrator\n"
    "├── scripts/\n"
    "│   ├── simulate_attack.c\n"
    "│   └── run_demo.sh\n"
    "├── examples/policy.yaml\n"
    "├── Makefile\n"
    "├── Dockerfile\n"
    "└── README.md", font_size=12)

# Build pipeline
add_shape(slide9, Inches(6.5), Inches(1.5), Inches(6.3), Inches(2.5), DARK_CARD, ACCENT_BLUE)
add_text_box(slide9, Inches(6.8), Inches(1.6), Inches(5.7), Inches(0.4),
             "Build Pipeline", font_size=18, color=ACCENT_BLUE, bold=True)
add_code_block(slide9, Inches(6.7), Inches(2.1), Inches(5.9), Inches(1.7),
    "# 1. Compile BPF (kernel program)\n"
    "clang -O2 -target bpf -c escape_detector.bpf.c\n"
    "llvm-strip -g escape_detector.bpf.o\n"
    "\n"
    "# 2. Compile detector (user-space)\n"
    "gcc -O2 detector.c -o detector -lbpf -lelf -lz\n"
    "\n"
    "# 3. Run\n"
    "sudo ./bin/detector", font_size=11)

# Technology stack
add_shape(slide9, Inches(6.5), Inches(4.3), Inches(6.3), Inches(2.6), DARK_CARD, ACCENT_GREEN)
add_text_box(slide9, Inches(6.8), Inches(4.4), Inches(5.7), Inches(0.4),
             "Technology Stack", font_size=18, color=ACCENT_GREEN, bold=True)

stack = [
    ("Kernel", "Linux 6.17, eBPF tracepoints, CO-RE"),
    ("BPF Toolchain", "clang/llvm (BPF target), llvm-strip"),
    ("Libraries", "libbpf, libelf, zlib"),
    ("Language", "C (kernel + user-space, ~850 lines total)"),
    ("Container", "Docker 28.x, ubuntu:22.04 images"),
    ("Platform", "Ubuntu 24.04 / aarch64"),
]
y = Inches(4.9)
for label, val in stack:
    txBox = slide9.shapes.add_textbox(Inches(6.8), y, Inches(5.7), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = f"{label}:  "
    r1.font.size = Pt(12); r1.font.color.rgb = ACCENT_GREEN; r1.font.bold = True; r1.font.name = "Segoe UI"
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(12); r2.font.color.rgb = LIGHT_GRAY; r2.font.name = "Segoe UI"
    y += Inches(0.32)

add_slide_number(slide9, 9)

# ══════════════════════════════════════════════════════════════════
#  SLIDE 10 — Summary & Conclusion
# ══════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide10)
add_shape(slide10, Inches(0), Inches(0), W, Pt(4), ACCENT_GREEN)

add_text_box(slide10, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "SUMMARY & CONCLUSION", font_size=32, color=WHITE, bold=True)
add_accent_line(slide10, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_GREEN)

# What We Built
add_shape(slide10, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.8), DARK_CARD, ACCENT_BLUE)
add_text_box(slide10, Inches(0.8), Inches(1.6), Inches(5.4), Inches(0.4),
             "What We Built", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide10, Inches(0.8), Inches(2.1), Inches(5.4), Inches(2.0), [
    "eBPF kernel program (191 lines) tracing 4 syscalls",
    "C user-space detector (658 lines) with 5 attack classifiers",
    "5 real attack scripts against live Docker containers",
    "CO-RE portable BPF: runs on any kernel 5.x+ without recompile",
    "Zero-dependency detection: no container agents needed",
], font_size=13, bullet_color=ACCENT_BLUE)

# Key Results
add_shape(slide10, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8), DARK_CARD, ACCENT_GREEN)
add_text_box(slide10, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4),
             "Key Results", font_size=20, color=ACCENT_GREEN, bold=True)

results = [
    ("5/5", "attack categories detected correctly"),
    ("0", "cross-triggers between detectors"),
    ("0", "false positives from host system noise"),
    ("< 1ms", "detection latency (kernel → alert)"),
    ("~850", "total lines of C code"),
]
y_r = Inches(2.15)
for metric, desc in results:
    txBox = slide10.shapes.add_textbox(Inches(7.1), y_r, Inches(5.4), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = f"{metric}  "
    r1.font.size = Pt(16); r1.font.color.rgb = ACCENT_GREEN; r1.font.bold = True; r1.font.name = "Consolas"
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(13); r2.font.color.rgb = LIGHT_GRAY; r2.font.name = "Segoe UI"
    y_r += Inches(0.38)

# Detection Coverage Table
add_shape(slide10, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), DARK_CARD, ACCENT_PURPLE)
add_text_box(slide10, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.4),
             "Detection Coverage Summary", font_size=18, color=ACCENT_PURPLE, bold=True)

# Table header
headers = ["Attack", "Docker Flag", "Syscall Traced", "Detection Rule", "Severity"]
col_x = [Inches(0.7), Inches(3.0), Inches(5.6), Inches(8.0), Inches(11.0)]
col_w = [Inches(2.2), Inches(2.5), Inches(2.3), Inches(2.9), Inches(1.5)]

y_t = Inches(5.15)
for i, h in enumerate(headers):
    add_text_box(slide10, col_x[i], y_t, col_w[i], Inches(0.3),
                 h, font_size=11, color=ACCENT_PURPLE, bold=True)

rows = [
    ["1. Docker Socket", "-v docker.sock", "connect(AF_UNIX)", "docker-socket-connect", "CRITICAL"],
    ["2. Privileged Escape", "--privileged", "mount() + openat()", "container-mount", "CRITICAL"],
    ["3. Cgroup Escape", "--privileged", "mount(cgroup)", "cgroup-mount", "CRITICAL"],
    ["4. Sensitive FS", "-v /:/hostfs:ro", "openat(hostfs/*)", "hostfs-credential-read", "HIGH"],
    ["5. Namespace Escape", "--privileged --pid=host", "setns() + openat(ns/)", "namespace-setns", "CRITICAL"],
]

y_t = Inches(5.45)
for row in rows:
    for i, cell in enumerate(row):
        c = ACCENT_GREEN if i == 4 and cell == "CRITICAL" else (ACCENT_ORANGE if i == 4 else LIGHT_GRAY)
        fn = "Consolas" if i in (2, 3) else "Segoe UI"
        add_text_box(slide10, col_x[i], y_t, col_w[i], Inches(0.25),
                     cell, font_size=10, color=c, font_name=fn)
    y_t += Inches(0.26)

add_slide_number(slide10, 10)

# ── Save ──────────────────────────────────────────────────────────
output_path = "/home/perry/6130 latest/Container_Runtime_Security_eBPF_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
